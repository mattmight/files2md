#!/usr/bin/env python3
"""
Enhanced File to Markdown Converter

This script combines multiple files into a single Markdown document, treating each file
as an "escaped attachment". It supports PDF, DOCX, PPTX (via python-pptx), plain text,
and XLSX files, with recursive directory searching and filename pattern matching.

Changes in this version
-----------------------
* **Hidden-file handling** – dotfiles are skipped by default.
  Pass `--include-hidden-files` to override.

Dependencies
------------
- pdfminer.six     : PDF text extraction
- pdf2image        : PDF → image conversion (for OCR fallback)
- pytesseract      : OCR on PDF images
- pypandoc         : DOCX → Markdown conversion
- python-pptx      : PPTX text extraction
- pandas + openpyxl: XLSX → Markdown tables
"""

import argparse
import os
import re
import sys

from pdf2image import convert_from_path
from pdfminer.high_level import extract_text
import pandas as pd
import pypandoc
import pytesseract

# Optional PPTX support
try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # handled gracefully at runtime

# --------------------------------------------------------------------------- #
#  Filename-matching presets
# --------------------------------------------------------------------------- #

CODE_FILE_PATTERN = (
    r'\.(py|js|java|c|cpp|h|hpp|cs|rb|go|rs|php|html|css|sql|sh|bash|ps1|rkt|hs|'
    r'scala|ml|elm|clj|ex|exs|erl|fs|fsx|lisp|scm|sml|swift|kt|kts|groovy|pl|pm|'
    r't|lua|jl|dart|d|nim|cr|r|R|asm|s|zig|v|ada|f90|f95|f03|f08|pas|cob|cobol|'
    r'vb|vba|vbs|tcl|hx|m|mm|ts|coffee|ls|cljc|cljs|raku|bf|md|txt)$'
    r'|^(Makefile|Dockerfile|Rakefile|Gemfile|Vagrantfile|CMakeLists\.txt)$'
)
DOCS_FILE_PATTERN = (
    r'\.(md|txt|rst|tex|rtf|odt|doc|docx|pdf|epub|csv|tsv|json|xml|yaml|yml|ini|'
    r'cfg|conf|log|pptx)$'
)

# --------------------------------------------------------------------------- #
#  Utility helpers
# --------------------------------------------------------------------------- #


def escape_backticks(text: str) -> str:
    """Wrap text in a back-tick fence long enough to avoid collisions."""
    ticks = max((len(m.group(0)) for m in re.finditer(r'`+', text)), default=0) + 1
    fence = '`' * max(3, ticks)
    return f"{fence}\n{text}\n{fence}"


def escape_xml_tags(text: str, filepath: str, base_tag: str = "file-attachment") -> tuple[str, str]:
    """Wrap text in XML-like tags, avoiding collisions with closing tags in content.

    Args:
        text: The content to wrap
        filepath: Path to the file (for the original attribute)
        base_tag: Base tag name (default: "file-attachment")

    Returns:
        Tuple of (wrapped_content, tag_used)
    """
    # Find a tag name that doesn't conflict with content
    tag_name = base_tag
    counter = 0

    while f"</{tag_name}>" in text:
        counter += 1
        tag_name = f"{base_tag}-{counter}"

    # Create the wrapped content with relative path attribute
    rel_path = os.path.relpath(filepath)
    opening_tag = f'<{tag_name} original="{rel_path}">'
    closing_tag = f'</{tag_name}>'

    wrapped_content = f"{opening_tag}\n{text}\n{closing_tag}"
    return wrapped_content, tag_name


# --------------------------------------------------------------------------- #
#  File-type converters
# --------------------------------------------------------------------------- #


def extract_text_from_pdf(path: str) -> str:
    """Extract text from PDF, OCR-fallback when necessary."""
    txt = extract_text(path)
    if txt.strip():
        return txt
    images = convert_from_path(path)
    return ''.join(pytesseract.image_to_string(img) for img in images)


def convert_docx_to_markdown(path: str) -> str:
    """DOCX → GitHub-flavoured Markdown via pandoc."""
    try:
        return pypandoc.convert_file(
            path, "gfm", format="docx",
            extra_args=["--wrap=none", "--standalone", "--markdown-headings=atx"]
        )
    except (RuntimeError, OSError) as e:
        return f"Failed to convert DOCX: {e}"


def convert_pptx_to_text(path: str) -> str:
    """Extract visible text from PPTX slides (requires python-pptx)."""
    if Presentation is None:
        return ("Failed to read PPTX: python-pptx not installed "
                "(pip install python-pptx)")
    try:
        prs = Presentation(path)
    except Exception as e:
        return f"Failed to open PPTX: {e}"

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = [f"# Slide {idx}"]
        for shape in slide.shapes:
            if getattr(shape, "text", "").strip():
                lines.append(shape.text.strip())
        if len(lines) == 1:
            lines.append("(No text on this slide)")
        slides.append("\n".join(lines))
    return "\n\n".join(slides)


def convert_xlsx_to_markdown(path: str) -> str:
    """Each sheet → GitHub-flavoured Markdown table."""
    try:
        sheets = pd.read_excel(path, sheet_name=None)
    except Exception as e:
        return f"Failed to read XLSX: {e}"

    parts = []
    for name, df in sheets.items():
        parts.append(f"### Sheet: {name}\n\n{df.to_markdown(index=False, tablefmt='github')}")
    return "\n\n".join(parts)


def convert_tsv_to_markdown(path: str) -> str:
    """Convert TSV file to GitHub-flavoured Markdown table."""
    try:
        df = pd.read_csv(path, sep='\t')
        return df.to_markdown(index=False, tablefmt='github')
    except Exception as e:
        return f"Failed to read TSV: {e}"


# --------------------------------------------------------------------------- #
#  Core processing
# --------------------------------------------------------------------------- #


def process_files(filepaths, name_regex=None, use_xml_tags=False):
    """Read, convert and wrap each file; return big Markdown string."""
    out, failures = [], []

    for path in filepaths:
        if name_regex and not re.search(name_regex, path):
            continue

        print(f"[files2md] ⌛ processing: {path}", file=sys.stderr)
        lower = path.lower()

        try:
            if lower.endswith(".pdf"):
                raw_content = extract_text_from_pdf(path)

            elif lower.endswith(".docx"):
                raw_content = convert_docx_to_markdown(path)

            elif lower.endswith(".pptx"):
                raw_content = convert_pptx_to_text(path)

            elif lower.endswith(".xlsx"):
                raw_content = convert_xlsx_to_markdown(path)

            elif lower.endswith(".tsv"):
                raw_content = convert_tsv_to_markdown(path)

            else:  # plain text or unknown
                try:
                    with open(path, encoding="utf-8") as fh:
                        raw_content = fh.read()
                except UnicodeDecodeError as e:
                    print(f"[files2md] ⚠️  UTF-8 failed for {path}: {e}", file=sys.stderr)
                    try:
                        with open(path, encoding="latin-1") as fh:
                            raw_content = fh.read()
                    except Exception as e2:
                        print(f"[files2md] 🚫 could not decode {path}: {e2}", file=sys.stderr)
                        failures.append(path)
                        continue

            # Apply appropriate escaping based on mode
            if use_xml_tags:
                content = raw_content
            else:
                # For non-markdown files, apply backtick escaping
                if lower.endswith((".xlsx", ".tsv")):
                    content = raw_content  # These are already markdown tables
                else:
                    content = escape_backticks(raw_content)

            end_comment = f"<!-- end of: {os.path.basename(path)} -->"
            if use_xml_tags:
                wrapped_content, tag_used = escape_xml_tags(content, path)
                out.extend([f"## Attached file: {path}", wrapped_content, end_comment])
            else:
                out.extend([f"<!-- file-attachment: {path} -->",
                            f"## Attached file: {path}", content, end_comment])

        except Exception as e:
            print(f"[files2md] 🚫 unexpected error on {path}: {e}", file=sys.stderr)
            failures.append(path)

    if failures:
        print("\n[files2md] SUMMARY — failed files:", file=sys.stderr)
        for f in failures:
            print(f"  • {f}", file=sys.stderr)

    return "\n\n".join(out)


# --------------------------------------------------------------------------- #
#  File discovery
# --------------------------------------------------------------------------- #


def is_hidden(path: str) -> bool:
    """True if the basename starts with a dot."""
    return os.path.basename(path).startswith(".")


def find_files(paths, *, recursive=False, include_hidden=False, name_regex=None):
    """Resolve all input paths to a flat list of filepaths."""
    results = []

    for path in paths:
        if not include_hidden and is_hidden(path):
            continue

        if os.path.isfile(path):
            results.append(path)

        elif os.path.isdir(path):
            if recursive:
                for root, dirs, files in os.walk(path):
                    if not include_hidden:
                        # prune hidden directories in-place
                        dirs[:] = [d for d in dirs if not d.startswith(".")]
                        files = [f for f in files if not f.startswith(".")]
                    for fname in files:
                        fpath = os.path.join(root, fname)
                        if include_hidden or not is_hidden(fpath):
                            if name_regex is None or re.search(name_regex, fpath):
                                results.append(fpath)
            else:
                for fname in os.listdir(path):
                    if not include_hidden and fname.startswith("."):
                        continue
                    fpath = os.path.join(path, fname)
                    if os.path.isfile(fpath) and (
                        name_regex is None or re.search(name_regex, fpath)
                    ):
                        results.append(fpath)

    return results


# --------------------------------------------------------------------------- #
#  CLI entry-point
# --------------------------------------------------------------------------- #


def main():
    try:
        pypandoc.get_pandoc_version()
    except OSError:
        print("Warning: pandoc not found – DOCX conversion disabled.", file=sys.stderr)

    ap = argparse.ArgumentParser(
        description="Combine files into a single Markdown document."
    )
    ap.add_argument("paths", nargs="+", help="Files or directories to process.")
    ap.add_argument(
        "-r", "--recursive", action="store_true",
        help="Recursively search directories."
    )
    hidden = ap.add_mutually_exclusive_group()
    hidden.add_argument(
        "--include-hidden-files", action="store_true", dest="include_hidden",
        help="Include dotfiles (hidden files)."
    )
    hidden.add_argument(
        "--exclude-hidden-files", action="store_false", dest="include_hidden",
        help="(Default) Skip dotfiles.", default=False
    )
    ap.add_argument("--name-regex", help="Regex to filter filenames.")
    ap.add_argument("--name-code", action="store_true",
                    help="Use preset regex for common code files.")
    ap.add_argument("--name-docs", action="store_true",
                    help="Use preset regex for common document files.")
    ap.add_argument("--xml-tags", action="store_true",
                    help="Use XML-like tags instead of backtick fences "
                         "(defaults to <file-attachment>).")
    args = ap.parse_args()

    if sum(map(bool, (args.name_regex, args.name_code, args.name_docs))) > 1:
        ap.error("Choose only one of --name-regex / --name-code / --name-docs.")

    regex = (
        args.name_regex if args.name_regex else
        CODE_FILE_PATTERN if args.name_code else
        DOCS_FILE_PATTERN if args.name_docs else
        None
    )

    files = find_files(
        args.paths,
        recursive=args.recursive,
        include_hidden=args.include_hidden,
        name_regex=regex,
    )
    print(process_files(files, name_regex=regex, use_xml_tags=args.xml_tags))


if __name__ == "__main__":
    main()
